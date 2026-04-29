from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== 设计系统 =====
BG = RGBColor(0xFA, 0xFA, 0xF8)
INK = RGBColor(0x1A, 0x1A, 0x1A)
INK_SEC = RGBColor(0x6B, 0x6B, 0x6B)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
LINE = RGBColor(0xE5, 0xE5, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def txt(txBox, text, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Microsoft YaHei"
    p.alignment = align

def footer(slide, page, total=19):
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.3))
    p = tx.text_frame.paragraphs[0]
    p.text = f"辟谣防卫战 · 毕业设计\t\t{page:02d} / {total:02d}"
    p.font.size = Pt(10); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"

def blue_bar(slide, x, y, h):
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.02), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

def tag(slide, label):
    tx = slide.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3))
    p = tx.text.frame.paragraphs[0]
    p.text = label; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"

def page_title(slide, title_text, page):
    h = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9))
    txt(h, title_text, 56, True, INK)

# ===== P01 封面 =====
s1 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s1)
lt = s1.shapes.add_shape(1, Inches(6.55), Inches(1.5), Inches(0.02), Inches(0.8)); lt.fill.solid(); lt.fill.fore_color.rgb = ACCENT; lt.line.fill.background()
t1 = s1.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.3), Inches(1)); txt(t1, "辟谣防卫战", 72, True, INK, PP_ALIGN.CENTER)
s1b = s1.shapes.add_textbox(Inches(1), Inches(3.6), Inches(11.3), Inches(0.6)); txt(s1b, "相亲相爱一家人", 24, False, INK_SEC, PP_ALIGN.CENTER)
dv = s1.shapes.add_shape(1, Inches(5.85), Inches(4.4), Inches(1.7), Inches(0.02)); dv.fill.solid(); dv.fill.fore_color.rgb = LINE; dv.line.fill.background()
auth = s1.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11.3), Inches(1.5))
tf = auth.text_frame; tf.word_wrap = True
for i, line in enumerate(["答辩人：XXX", "设计学 · 视觉传达设计", "指导教师：XXX 教授", "2026 年 6 月"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line; p.font.size = Pt(14 if i == 0 else 12); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
lb = s1.shapes.add_shape(1, Inches(6.55), Inches(6.3), Inches(0.02), Inches(0.8)); lb.fill.solid(); lb.fill.fore_color.rgb = ACCENT; lb.line.fill.background()
yr = s1.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11.3), Inches(0.4)); txt(yr, "THESIS DEFENSE · 2026", 10, False, INK_SEC, PP_ALIGN.CENTER)
footer(s1, 1)

# ===== P02 目录 =====
s2 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s2)
pt2 = s2.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(3), Inches(0.4)); txt(pt2, "CONTENTS", 11, False, INK_SEC)
chaps = [
    ("01", "研究背景", "Research Background", ["选题背景 · 问题提出 · 研究意义", "3 页"]),
    ("02", "设计策略", "Design Strategy", ["用户调研 · 竞品分析 · 核心洞察", "设计概念 · 系统架构 · 信息结构", "6 页"]),
    ("03", "成果展示", "Output & Showcase", ["方案推导 · 主要功能 · 核心亮点", "视觉规范 · 成果展示 · 展示方式", "总结与展望 · 致谢", "10 页"]),
]
for i, (num, tc, sub, lines) in enumerate(chaps):
    x = 0.8 + i * 4.2
    blue_bar(s2, x - 0.01, 1.8, 4.2)
    nb = s2.shapes.add_textbox(Inches(x), Inches(1.8), Inches(3.8), Inches(1)); txt(nb, num, 64, True, ACCENT)
    tb = s2.shapes.add_textbox(Inches(x), Inches(2.9), Inches(3.8), Inches(0.5)); txt(tb, tc, 24, True, INK)
    sb = s2.shapes.add_textbox(Inches(x), Inches(3.35), Inches(3.8), Inches(0.4)); txt(sb, sub, 10, False, INK_SEC)
    pb = s2.shapes.add_textbox(Inches(x), Inches(4.9), Inches(3.8), Inches(1.5))
    tf2 = pb.text_frame; tf2.word_wrap = True
    for j, line in enumerate(lines):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.text = line; p.font.size = Pt(11); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"
footer(s2, 2)

# ===== P03 选题背景 =====
s3 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s3)
tag3 = s3.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag3.text_frame.paragraphs[0]; p.text = "研究背景 01"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h3 = s3.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h3, "选题背景", 56, True, INK)
cb3 = s3.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(5.5), Inches(2))
tf3 = cb3.text_frame; tf3.word_wrap = True
for i, line in enumerate(["家族微信群已成为谣言传播的主要渠道。", "老年人因缺乏辨别能力，往往成为谣言的二次传播者。", "传统辟谣方式说教性强、接受度低。"]):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    p.text = line; p.font.size = Pt(18); p.font.color.rgb = INK; p.font.name = "Microsoft YaHei"; p.line_spacing = 2
dcs = [("68%", "中老年人曾在家族群转发谣言"), ("92%", "受访年轻人表示曾被谣言困扰"), ("3次", "平均每人每日接触谣言次数")]
for i, (n, d) in enumerate(dcs):
    y = 2.5 + i * 1.3
    blue_bar(s3, 7.2, y, 1.0)
    nb = s3.shapes.add_textbox(Inches(7.4), Inches(y - 0.05), Inches(2), Inches(0.55)); txt(nb, n, 36, True, ACCENT)
    db = s3.shapes.add_textbox(Inches(7.4), Inches(y + 0.55), Inches(5), Inches(0.6)); txt(db, d, 13, False, INK_SEC)
footer(s3, 3)

# ===== P04 问题提出 =====
s4 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s4)
tag4 = s4.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag4.text_frame.paragraphs[0]; p.text = "研究背景 02"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h4, "问题提出", 56, True, INK)
lims = [
    ("01", "科学理性不足", "官方媒体辟谣文章往往充满专业术语和数据，普通老年人难以理解消化。缺乏通俗易懂的表达方式，导致辟谣效果有限。"),
    ("02", "情感共鸣缺失", "传统辟谣方式说教感强，单方面输出结论，忽略了与受众的情感连接。难以触及老年人的情感需求，认同感低。"),
    ("03", "行动转化困难", "缺乏有效的行动引导，用户知道是谣言但不知道如何应对。缺少互动参与感，无法转化为长期的防御能力。"),
]
for i, (n, t, d) in enumerate(lims):
    x = 0.8 + i * 4.1; y = 2.3
    card = s4.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(4))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    top = s4.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(0.04))
    top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()
    nb = s4.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(3.4), Inches(0.8)); txt(nb, n, 48, True, ACCENT)
    tb = s4.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.2), Inches(3.4), Inches(0.5)); txt(tb, t, 20, True, INK)
    db = s4.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.8), Inches(3.4), Inches(2))
    tf4 = db.text_frame; tf4.word_wrap = True
    p = tf4.paragraphs[0]; p.text = d; p.font.size = Pt(13); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"; p.line_spacing = 1.8
footer(s4, 4)

# ===== P05 研究意义 =====
s5 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s5)
tag5 = s5.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag5.text_frame.paragraphs[0]; p.text = "研究背景 03"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h5, "研究意义", 56, True, INK)
main5 = s5.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5), Inches(1.5))
tf5 = main5.text_frame; tf5.word_wrap = True
for i, line in enumerate(["本研究旨在通过游戏化设计，", "创造一种轻量化、情感化、", "参与感强的辟谣体验。"]):
    p = tf5.paragraphs[0] if i == 0 else tf5.add_paragraph()
    p.text = line; p.font.size = Pt(20); p.font.color.rgb = INK; p.font.name = "Microsoft YaHei"; p.line_spacing = 1.8
goals = [
    ("轻量化体验", "无需阅读长文，5分钟一局即可完成辟谣任务"),
    ("情感化沟通", "通过MBTI人格系统建立共情，理解不同角色的谣言偏好"),
    ("游戏化学习", "在娱乐中自然建立谣言识别能力，告别被动说教"),
]
for i, (t, d) in enumerate(goals):
    x = 0.8 + i * 4.1; y = 4.2
    card = s5.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(2))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    blue_bar(s5, x, y, 2)
    tb = s5.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.2), Inches(3.3), Inches(0.5)); txt(tb, t, 18, True, INK)
    db = s5.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.7), Inches(3.3), Inches(1)); tfd = db.text_frame; tfd.word_wrap = True
    p = tfd.paragraphs[0]; p.text = d; p.font.size = Pt(13); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"
footer(s5, 5)

# ===== P06 用户调研 =====
s6 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s6)
tag6 = s6.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag6.text_frame.paragraphs[0]; p.text = "设计策略 01"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h6, "用户调研", 56, True, INK)
# 左侧用户画像
card6 = s6.shapes.add_shape(1, Inches(0.8), Inches(2.3), Inches(5), Inches(3.8))
card6.fill.solid(); card6.fill.fore_color.rgb = WHITE; card6.line.fill.background()
pt6 = s6.shapes.add_textbox(Inches(1), Inches(2.5), Inches(4.6), Inches(0.4)); txt(pt6, "目标用户画像", 14, True, INK)
mbti6 = s6.shapes.add_textbox(Inches(1), Inches(2.95), Inches(4.6), Inches(0.35)); txt(mbti6, "ENFP · 焦虑传播型", 12, False, ACCENT)
desc6 = s6.shapes.add_textbox(Inches(1), Inches(3.35), Inches(4.6), Inches(2.5))
tfd6 = desc6.text_frame; tfd6.word_wrap = True
p = tfd6.paragraphs[0]; p.text = "容易被情感化内容打动，热衷转发「好心」提醒。收到「养生忠告」会立刻转发到家族群，不验证真伪。"; p.font.size = Pt(14); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"; p.line_spacing = 1.8
# 右侧洞察
insights = [
    ("洞察 01", "不同MBTI人格对谣言类型有明确偏好，ISTJ倾向相信食品安全类，ESFJ易被情感谣言打动"),
    ("洞察 02", "用户不拒绝辟谣内容，但拒绝「被教育」的感觉，需要平等的沟通姿态"),
    ("洞察 03", "游戏化的互动方式参与度比静态内容高3倍，动机转化效果显著"),
    ("洞察 04", "立即反馈机制比延迟说教更有效"),
]
for i, (t, d) in enumerate(insights):
    x = 6.3; y = 2.3 + i * 1.05
    card_i = s6.shapes.add_shape(1, Inches(x), Inches(y), Inches(6.2), Inches(0.9))
    card_i.fill.solid(); card_i.fill.fore_color.rgb = WHITE; card_i.line.fill.background()
    tb = s6.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(5.8), Inches(0.35)); txt(tb, t, 12, True, ACCENT)
    db = s6.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.45), Inches(5.8), Inches(0.5)); txt(db, d, 13, False, INK_SEC)
footer(s6, 6)

# ===== P07 竞品分析 =====
s7 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s7)
tag7 = s7.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag7.text_frame.paragraphs[0]; p.text = "设计策略 02"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h7 = s7.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h7, "竞品分析", 56, True, INK)
pros = [("辟谣公众号", "内容专业", "缺乏互动、说教感强"), ("辟谣小程序", "轻量便捷", "功能单一、深度不足"), ("长辈表情包", "接受度高", "深度有限、无法纠错")]
for i, (name, pro, con) in enumerate(pros):
    x = 0.8 + i * 4.1; y = 2.3
    card = s7.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(3.8))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    nb = s7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(3.4), Inches(0.5)); txt(nb, name, 20, True, INK)
    pb = s7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9), Inches(3.4), Inches(0.3)); txt(pb, "优势：" + pro, 13, False, RGBColor(0x05, 0x96, 0x69))
    cb = s7.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.25), Inches(3.4), Inches(0.3)); txt(cb, "局限：" + con, 13, False, RGBColor(0xDC, 0x26, 0x26))
footer(s7, 7)

# ===== P08 核心洞察 =====
s8 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s8)
tag8 = s8.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag8.text_frame.paragraphs[0]; p.text = "设计策略 03"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h8 = s8.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h8, "核心洞察", 56, True, INK)
# 大引用
quote8 = s8.shapes.add_shape(1, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.2))
quote8.fill.solid(); quote8.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF); quote8.line.fill.background()
qt = s8.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11), Inches(1)); txt(qt, "「不同人格的辟谣接受度差异，是设计游戏化体验的核心机会。」", 22, False, INK)
insights8 = [
    ("01", "MBTI性格决定信谣类型", "不同人格对谣言有明确偏好，理解人格才能对症下药"),
    ("02", "游戏化提升参与动机", "竞争与奖励机制能将被动辟谣转化为主动学习"),
    ("03", "情感共鸣胜于理性说教", "平等对话比教导更有效，温柔沟通比强硬指出更易接受"),
]
for i, (n, t, d) in enumerate(insights8):
    x = 0.8 + i * 3.9; y = 3.9
    blue_bar(s8, x, y, 2.2)
    nb = s8.shapes.add_textbox(Inches(x + 0.2), Inches(y), Inches(3.5), Inches(0.6)); txt(nb, n, 32, True, ACCENT)
    tb = s8.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.65), Inches(3.5), Inches(0.5)); txt(tb, t, 16, True, INK)
    db = s8.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.2), Inches(3.5), Inches(1)); txt(db, d, 13, False, INK_SEC)
footer(s8, 8)

# ===== P09 设计概念 =====
s9 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s9)
tag9 = s9.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag9.text_frame.paragraphs[0]; p.text = "设计策略 04"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h9, "设计概念", 56, True, INK)
concept9 = s9.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(5.5), Inches(2))
tf9 = concept9.text_frame; tf9.word_wrap = True
for i, line in enumerate(["通过 MBTI 性格系统理解不同角色的谣言偏好，", "借助卡牌游戏机制做出辟谣选择，", "在双 HP 博弈中平衡家庭和谐与谣言污染。"]):
    p = tf9.paragraphs[0] if i == 0 else tf9.add_paragraph()
    p.text = line; p.font.size = Pt(18); p.font.color.rgb = INK; p.font.name = "Microsoft YaHei"; p.line_spacing = 2
# 右侧流程
flow_items = [("输入", "家族群谣言"), ("人格选择", "理解不同人格的谣言偏好"), ("卡牌辟谣", "6选1，即时反馈"), ("HP变化", "和谐值与污染值的博弈"), ("输出", "辟谣能力提升")]
for i, (t, d) in enumerate(flow_items):
    y = 2.4 + i * 0.75
    if i == 0 or i == 4:
        box = s9.shapes.add_shape(1, Inches(7), Inches(y), Inches(5), Inches(0.6))
        box.fill.solid(); box.fill.fore_color.rgb = INK; box.line.fill.background()
        tb = s9.shapes.add_textbox(Inches(7.2), Inches(y + 0.1), Inches(4.6), Inches(0.5)); txt(tb, t + " · " + d, 16, True, WHITE)
    else:
        box = s9.shapes.add_shape(1, Inches(7), Inches(y), Inches(5), Inches(0.6))
        box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.fill.background()
        tb = s9.shapes.add_textbox(Inches(7.2), Inches(y + 0.1), Inches(4.6), Inches(0.5)); txt(tb, t + " · " + d, 16, False, INK)
footer(s9, 9)

# ===== P10 系统架构 =====
s10 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s10)
tag10 = s10.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag10.text_frame.paragraphs[0]; p.text = "设计策略 05"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h10 = s10.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h10, "系统架构", 56, True, INK)
# 双HP说明
hp_items = [
    ("家庭和谐值", "100", "玩家需要保护的目标，归零则输", RGBColor(0x05, 0x96, 0x69)),
    ("谣言污染值", "100", "玩家需要消除的目标，归零则赢", RGBColor(0xDC, 0x26, 0x26)),
]
for i, (name, val, desc, col) in enumerate(hp_items):
    y = 2.4 + i * 1.5
    bar = s10.shapes.add_shape(1, Inches(0.8), Inches(y), Inches(0.04), Inches(1.1)); bar.fill.solid(); bar.fill.fore_color.rgb = col; bar.line.fill.background()
    nb = s10.shapes.add_textbox(Inches(1.1), Inches(y), Inches(2), Inches(0.5)); txt(nb, name, 22, True, INK)
    vb = s10.shapes.add_textbox(Inches(1.1), Inches(y + 0.5), Inches(1.5), Inches(0.5)); txt(vb, val, 36, True, col)
    db = s10.shapes.add_textbox(Inches(2.8), Inches(y + 0.5), Inches(4), Inches(0.5)); txt(db, desc, 14, False, INK_SEC)
# 回合流程
flow10 = [("1", "展示谣言"), ("2", "选择卡牌"), ("3", "计算HP"), ("4", "判定胜负")]
for i, (n, t) in enumerate(flow10):
    x = 7.2 + i * 1.5; y = 2.8
    circle = s10.shapes.add_shape(9, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
    nt = s10.shapes.add_textbox(Inches(x), Inches(y + 0.05), Inches(0.5), Inches(0.4)); txt(nt, n, 18, True, WHITE, PP_ALIGN.CENTER)
    tb = s10.shapes.add_textbox(Inches(x - 0.3), Inches(y + 0.6), Inches(1.1), Inches(0.5)); txt(tb, t, 12, False, INK_SEC, PP_ALIGN.CENTER)
    if i < 3:
        arrow = s10.shapes.add_textbox(Inches(x + 0.55), Inches(y + 0.1), Inches(0.5), Inches(0.3)); txt(arrow, "→", 18, False, INK_SEC)
footer(s10, 10)

# ===== P11 信息结构 =====
s11 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s11)
tag11 = s11.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag11.text_frame.paragraphs[0]; p.text = "设计策略 06"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h11 = s11.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h11, "信息结构", 56, True, INK)
cats = [
    ("健康养生", "14条", "微波炉致癌、酸碱体质、食物相克"),
    ("食品安全", "14条", "转基因、地沟油、隔夜菜"),
    ("社会民生", "14条", "偷孩子谣言、停水通知"),
    ("阴谋论", "13条", "专家被收买、资本阴谋"),
    ("科技类", "13条", "5G辐射、手机辐射、AI取代人类"),
]
for i, (name, count, example) in enumerate(cats):
    x = 0.8 + i * 2.35; y = 2.3
    card = s11.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.15), Inches(3.8))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    top = s11.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.15), Inches(0.04))
    top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()
    nb = s11.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.3), Inches(1.85), Inches(0.8)); txt(nb, count, 36, True, ACCENT)
    tb = s11.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.2), Inches(1.85), Inches(0.5)); txt(tb, name, 18, True, INK)
    eb = s11.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.8), Inches(1.85), Inches(2)); tfd = eb.text_frame; tfd.word_wrap = True
    p = tfd.paragraphs[0]; p.text = example; p.font.size = Pt(12); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"
footer(s11, 11)

# ===== P12 方案推导 =====
s12 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s12)
tag12 = s12.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag12.text_frame.paragraphs[0]; p.text = "成果展示 01"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h12 = s12.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h12, "方案推导", 56, True, INK)
steps12 = [("问题提出", "家族群谣言传播失控"), ("核心洞察", "MBTI人格影响辟谣接受度"), ("设计假设", "游戏化能提升辟谣动机"), ("方案产出", "MBTI+卡牌+双HP系统")]
for i, (t, d) in enumerate(steps12):
    x = 0.8 + i * 2.9; y = 2.6
    box = s12.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.6), Inches(3.2))
    box.fill.solid(); box.fill.fore_color.rgb = WHITE; box.line.fill.background()
    nb = s12.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(2.2), Inches(0.8)); txt(nb, f"0{i+1}", 48, True, ACCENT)
    tb = s12.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.3), Inches(2.2), Inches(0.5)); txt(tb, t, 18, True, INK)
    db = s12.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.9), Inches(2.2), Inches(1)); tfd = db.text_frame; tfd.word_wrap = True
    p = tfd.paragraphs[0]; p.text = d; p.font.size = Pt(13); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"
    if i < 3:
        arrow = s12.shapes.add_textbox(Inches(x + 2.65), Inches(y + 1.5), Inches(0.3), Inches(0.4)); txt(arrow, "→", 24, False, ACCENT)
footer(s12, 12)

# ===== P13 主要功能 =====
s13 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s13)
tag13 = s13.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag13.text_frame.paragraphs[0]; p.text = "成果展示 02"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h13 = s13.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h13, "主要功能", 56, True, INK)
funcs = [
    ("人格选择", "6种MBTI角色，各有独特的谣言抗性和辟谣难度"),
    ("卡牌对战", "每回合6张卡牌，2正确+4错误，即时反馈"),
    ("双HP系统", "家庭和谐值vs谣言污染值，策略两难"),
    ("科普弹窗", "选错时显示解释，教育性强"),
]
for i, (t, d) in enumerate(funcs):
    y = 2.3 + i * 1.05
    blue_bar(s13, 0.8, y, 0.9)
    tb = s13.shapes.add_textbox(Inches(1.1), Inches(y + 0.1), Inches(3), Inches(0.4)); txt(tb, t, 18, True, INK)
    db = s13.shapes.add_textbox(Inches(1.1), Inches(y + 0.5), Inches(10), Inches(0.5)); txt(db, d, 14, False, INK_SEC)
footer(s13, 13)

# ===== P14 核心亮点 =====
s14 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s14)
tag14 = s14.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag14.text_frame.paragraphs[0]; p.text = "成果展示 03"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h14 = s14.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h14, "核心亮点", 56, True, INK)
highlights = [
    ("68+", "真实谣言", "覆盖5大类别", "内容深度"),
    ("6", "MBTI角色", "各有信谣弱点", "角色设计"),
    ("2", "HP双系统", "策略两难博弈", "机制创新"),
    ("0", "网络依赖", "完全离线可用", "离线可用"),
]
for i, (n, t, d, tag_l) in enumerate(highlights):
    x = 0.8 + i * 3.1; y = 2.3
    card = s14.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(4))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    top = s14.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(0.04))
    top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()
    nb = s14.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.3), Inches(2.6), Inches(1)); txt(nb, n, 64, True, INK)
    tb = s14.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.4), Inches(2.6), Inches(0.5)); txt(tb, t, 20, True, INK)
    db = s14.shapes.add_textbox(Inches(x + 0.15), Inches(y + 2), Inches(2.6), Inches(0.8)); tfd = db.text_frame; tfd.word_wrap = True
    p = tfd.paragraphs[0]; p.text = d; p.font.size = Pt(13); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"
    tg = s14.shapes.add_textbox(Inches(x + 0.15), Inches(y + 3.3), Inches(2.6), Inches(0.4)); txt(tg, tag_l, 11, False, ACCENT)
footer(s14, 14)

# ===== P15 视觉规范 =====
s15 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s15)
tag15 = s15.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag15.text_frame.paragraphs[0]; p.text = "成果展示 04"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h15 = s15.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h15, "视觉规范", 56, True, INK)
# 色彩
colors15 = [("#FAFAF8", "暖白背景"), ("#1A1A1A", "深灰文字"), ("#2563EB", "蓝色强调"), ("#E5E5E5", "分割线")]
for i, (c, d) in enumerate(colors15):
    x = 0.8 + i * 2.9
    swatch = s15.shapes.add_shape(1, Inches(x), Inches(2.4), Inches(0.8), Inches(0.8))
    swatch.fill.solid(); swatch.fill.fore_color.rgb = RGBColor(int(c[1:3],16), int(c[3:5],16), int(c[5:7],16)); swatch.line.fill.background()
    cb = s15.shapes.add_textbox(Inches(x), Inches(3.3), Inches(2.5), Inches(0.3)); txt(cb, c, 12, True, INK)
    db = s15.shapes.add_textbox(Inches(x), Inches(3.6), Inches(2.5), Inches(0.3)); txt(db, d, 11, False, INK_SEC)
# 字体层级
type_items = [("标题", "72px / 56px", "加粗"), ("副标题", "24px", "常规"), ("正文", "18px / 14px", "常规"), ("标签", "10px / 11px", "中等")]
for i, (t, s, w) in enumerate(type_items):
    x = 0.8 + i * 3
    tb = s15.shapes.add_textbox(Inches(x), Inches(4.4), Inches(2.8), Inches(0.4)); txt(tb, t, 16, True, INK)
    sb = s15.shapes.add_textbox(Inches(x), Inches(4.85), Inches(2.8), Inches(0.3)); txt(sb, s, 13, False, INK_SEC)
    wb = s15.shapes.add_textbox(Inches(x), Inches(5.15), Inches(2.8), Inches(0.3)); txt(wb, w, 11, False, ACCENT)
footer(s15, 15)

# ===== P16 成果展示 =====
s16 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s16)
tag16 = s16.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag16.text_frame.paragraphs[0]; p.text = "成果展示 05"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h16 = s16.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h16, "成果展示", 56, True, INK)
scenes = [("在线体验", "Netlify托管"), ("本地离线包", "123KB一键启动"), ("APK安装包", "手机端独立运行")]
for i, (t, d) in enumerate(scenes):
    x = 0.8 + i * 4.1; y = 2.5
    card = s16.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(3))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE; card.line.fill.background()
    top = s16.shapes.add_shape(1, Inches(x), Inches(y), Inches(3.8), Inches(0.04))
    top.fill.solid(); top.fill.fore_color.rgb = ACCENT; top.line.fill.background()
    tb = s16.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.4), Inches(3.4), Inches(0.5)); txt(tb, t, 22, True, INK)
    db = s16.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1), Inches(3.4), Inches(1.5)); tfd = db.text_frame; tfd.word_wrap = True
    p = tfd.paragraphs[0]; p.text = d; p.font.size = Pt(16); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"
footer(s16, 16)

# ===== P17 展示方式 =====
s17 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s17)
tag17 = s17.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag17.text_frame.paragraphs[0]; p.text = "成果展示 06"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h17 = s17.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h17, "展示方式", 56, True, INK)
methods = [
    ("在线版本", "https://fabulous-dasik-41be94.netlify.app/", "扫码即可体验"),
    ("本地离线包", "辟谣防卫战-展览体验版.zip", "解压双击start.bat"),
    ("APK安装包", "手机直接安装", "完全独立运行"),
]
for i, (t, url, desc) in enumerate(methods):
    y = 2.3 + i * 1.3
    blue_bar(s17, 0.8, y, 1.1)
    tb = s17.shapes.add_textbox(Inches(1.1), Inches(y + 0.1), Inches(3), Inches(0.4)); txt(tb, t, 18, True, INK)
    ub = s17.shapes.add_textbox(Inches(1.1), Inches(y + 0.5), Inches(10), Inches(0.3)); txt(ub, url, 12, False, ACCENT)
    db = s17.shapes.add_textbox(Inches(1.1), Inches(y + 0.8), Inches(10), Inches(0.3)); txt(db, desc, 12, False, INK_SEC)
footer(s17, 17)

# ===== P18 总结与展望 =====
s18 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s18)
tag18 = s18.shapes.add_textbox(Inches(11.2), Inches(0.8), Inches(1.5), Inches(0.3)); p = tag18.text_frame.paragraphs[0]; p.text = "成果展示 07"; p.font.size = Pt(10); p.font.color.rgb = ACCENT; p.font.name = "Microsoft YaHei"
h18 = s18.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8), Inches(0.9)); txt(h18, "总结与展望", 56, True, INK)
# 左：成果
card18a = s18.shapes.add_shape(1, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4))
card18a.fill.solid(); card18a.fill.fore_color.rgb = WHITE; card18a.line.fill.background()
blue_bar(s18, 0.8, 2.3, 4)
ta = s18.shapes.add_textbox(Inches(1.1), Inches(2.5), Inches(5), Inches(0.5)); txt(ta, "已完成", 20, True, INK)
items18a = ["MBTI人格系统设计", "68条真实谣言内容", "双HP博弈机制", "8-bit像素风格", "错误科普弹窗"]
for i, item in enumerate(items18a):
    y = 3.1 + i * 0.6
    db = s18.shapes.add_textbox(Inches(1.1), Inches(y), Inches(5), Inches(0.5)); txt(db, "✓ " + item, 14, False, INK_SEC)
# 右：展望
card18b = s18.shapes.add_shape(1, Inches(6.8), Inches(2.3), Inches(5.5), Inches(4))
card18b.fill.solid(); card18b.fill.fore_color.rgb = WHITE; card18b.line.fill.background()
bb = s18.shapes.add_shape(1, Inches(6.8), Inches(2.3), Inches(0.04), Inches(4)); bb.fill.solid(); bb.fill.fore_color.rgb = ACCENT; bb.line.fill.background()
tb18 = s18.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5), Inches(0.5)); txt(tb18, "未来方向", 20, True, INK)
items18b = ["扩充谣言库至100+条", "新增更多MBTI角色变体", "支持多人联机对战", "导出适配更多平台"]
for i, item in enumerate(items18b):
    y = 3.1 + i * 0.6
    db = s18.shapes.add_textbox(Inches(7.1), Inches(y), Inches(5), Inches(0.5)); txt(db, "→ " + item, 14, False, INK_SEC)
footer(s18, 18)

# ===== P19 致谢 =====
s19 = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s19)
lt19 = s19.shapes.add_shape(1, Inches(6.55), Inches(1.5), Inches(0.02), Inches(0.8)); lt19.fill.solid(); lt19.fill.fore_color.rgb = ACCENT; lt19.line.fill.background()
t19 = s19.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1)); txt(t19, "感谢聆听", 64, True, INK, PP_ALIGN.CENTER)
s19b = s19.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.3), Inches(0.5)); txt(s19b, "THANK YOU FOR YOUR ATTENTION", 16, False, INK_SEC, PP_ALIGN.CENTER)
dv19 = s19.shapes.add_shape(1, Inches(5.85), Inches(4.7), Inches(1.7), Inches(0.02)); dv19.fill.solid(); dv19.fill.fore_color.rgb = LINE; dv19.line.fill.background()
info19 = s19.shapes.add_textbox(Inches(1), Inches(5.1), Inches(11.3), Inches(0.8))
tf19 = info19.text_frame; tf19.word_wrap = True
for i, line in enumerate(["辟谣防卫战 · 毕业设计", "设计学 · 视觉传达设计"]):
    p = tf19.paragraphs[0] if i == 0 else tf19.add_paragraph()
    p.text = line; p.font.size = Pt(14); p.font.color.rgb = INK_SEC; p.font.name = "Microsoft YaHei"; p.alignment = PP_ALIGN.CENTER
lb19 = s19.shapes.add_shape(1, Inches(6.55), Inches(6.3), Inches(0.02), Inches(0.8)); lb19.fill.solid(); lb19.fill.fore_color.rgb = ACCENT; lb19.line.fill.background()
yr19 = s19.shapes.add_textbox(Inches(1), Inches(6.7), Inches(11.3), Inches(0.4)); txt(yr19, "2026", 12, False, INK_SEC, PP_ALIGN.CENTER)
footer(s19, 19)

# ===== 保存 =====
out = "E:/projects/BISHE/game/毕业设计答辩PPT_Final.pptx"
prs.save(out)
print(f"Done: {out}")